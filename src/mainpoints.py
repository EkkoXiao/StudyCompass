import PyPDF2
import fitz  # PyMuPDF
import os
from openai import OpenAI
import re
import hashlib # PyMuPDF
from langchain.document_loaders import UnstructuredPowerPointLoader, PyPDFLoader
from langchain.chat_models import ChatOpenAI
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import Chroma
from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
from langchain.chains.retrieval_qa.base import RetrievalQA
from langchain.document_loaders import TextLoader
import tempfile
from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE  
from PIL import Image  
from typing import Tuple
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import ast
from openai import OpenAI
import re
import docx
import docx2txt
from utils import *
from guideline import * 
import aspose.words as aw
from docx import Document

def process_folder(folder_path):
    """
    Processes each file in the given folder based on the file type (PPT, Word, or PDF).

    Args:
    folder_path (str): The path of the folder containing the files to process.
    """
    i = 0
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        # Ensure it's a file and not a directory
        if os.path.isfile(file_path):
            if filename.endswith('.pptx'):
                metadata = extract_metadata_from_pptx_langchain(file_path)
                extract_images_from_pptx(file_path,f"image{i}")
                key_word = key_word_gen(metadata)
                main_knowledge = ppt_knowledge_points_gen(metadata)
            elif filename.endswith('.docx'):
                metadata = extract_text_from_docx(file_path)
                extract_images_from_docx(file_path,f"image{i}")
                key_word = key_word_gen(metadata)
                main_knowledge = knowledge_points_gen(metadata)
            elif filename.endswith('.pdf'):
                metadata = extract_text_with_page_numbers_pdf(file_path)
                extract_images_from_pdf(file_path,f"image{i}")
                key_word = key_word_gen(metadata)
                main_knowledge = knowledge_points_gen(metadata)
            i += 1

# Example usage of process_folder
# process_folder('/path/to/your/folder')


shape_type_mapping = {
    1: "auto_shape",       # 自动形状
    2: "callout",          # 标注
    3: "chart",            # 图表
    4: "comment",          # 评论
    5: "free",         # 任意形状
    6: "group",            # 组合形状
    7: "embed", # 嵌入的OLE对象
    8: "control",     # 表单控件
    9: "LINE",             # 线条
    10: "LINKED_OLE_OBJECT", # 链接的OLE对象
    11: "LINKED_PICTURE",  # 链接的图片
    12: "OLE_CONTROL_OBJECT", # OLE控件对象
    13: "PICTURE",         # 图片
    14: "PLACEHOLDER",     # 占位符
    15: "text_box",        # 文本框
    16: "MEDIA",           # 媒体
    17: "table",           # 表格
    18: "TEXT_EFFECT",     # 文本效果
    19: "WEB_VIDEO",       # 网络视频
    20: "INK",             # 墨迹
    21: "INK_COMMENT",     # 墨迹评论
    22: "SMART_ART",       # 智能艺术
    23: "CONTENT_APP"      # 内容应用
}

def extract_metadata_from_pptx(path: str):
    presentation = Presentation(path)
    extracted_data = []

    for page, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                shape_type = shape.shape_type  # 获取形状类型
                content_type = "other"

 
                if shape_type == 14 and shape.placeholder_format.idx == 0:
                    content_type = "Title"
                elif shape_type == 14 and shape.placeholder_format.idx == 1:
                    content_type = "body"
                elif shape_type <= 23 :
                    content_type = shape_type_mapping[shape_type]

                data = {
                    "page": page, # 所在页数
                    "type": content_type, ## 类型
                    "content": text  ## 内容
                }
                extracted_data.append(data)

    return extracted_data


def extract_metadata_from_pptx_langchain(path: str):
    loader = UnstructuredPowerPointLoader(path,mode="elements")
    slides = loader.load()
    extracted_data = []
    page = 1
    for slide in slides:
        if slide.metadata.get('page_number') != None :
            page = slide.metadata.get('page_number')
        data = {
                        "page": page, # 所在页数
                        "type": slide.metadata.get('category'), ## 类型
                        "content": slide.page_content ## 内容
            }
        extracted_data.append(data)
    return extracted_data


def extract_images_from_pptx(path: str, output_folder: str):
    # 加载 PowerPoint 文件
    presentation = Presentation(path)
    slides = presentation.slides

    # 确保输出文件夹存在
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # 遍历每个幻灯片
    for slide_number, slide in enumerate(slides, start=1):
        image_count_on_slide = 0  # 重置每张幻灯片上的图片计数器

        # 遍历幻灯片中的每个形状
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 提取图片
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext

                # 保存图片
                image_filename = f"{slide_number}-{image_count_on_slide + 1}.{image_format}"
                image_path = os.path.join(output_folder, image_filename)
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)

                image_count_on_slide += 1

    print(f"Extracted images from {len(slides)} slides.")



def extract_images_from_pdf(pdf_path, output_folder):
    """
    Extracts images from a PDF and saves them in the specified output folder with the naming format 'xx-xx',
    where the first 'xx' is the page number and the second 'xx' is the image number on that page, avoiding duplicates.
    
    Args:
    pdf_path (str): Path to the PDF file.
    output_folder (str): Folder where extracted images will be saved.
    """
    # Ensure output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open the PDF
    doc = fitz.open(pdf_path)

    # Set to store image hashes to avoid duplicates
    seen_images = set()

    # Extract images
    for i in range(len(doc)):
        image_count = 0  # Initialize a counter for images on each page
        for img in doc.get_page_images(i):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]

            # Compute the hash of the image
            image_hash = hashlib.md5(image_bytes).hexdigest()

            # Skip the image if it's already been processed
            if image_hash in seen_images:
                continue
            seen_images.add(image_hash)

            # Increment the image count for this page
            image_count += 1

            # Save the image with the new naming format
            image_filename = f"{i+1}-{image_count}.{base_image['ext']}"
            with open(os.path.join(output_folder, image_filename), "wb") as f:
                f.write(image_bytes)

    # Close the document
    doc.close()



def extract_text_with_page_numbers_pdf(pdf_path):
    """
    Extracts text from a PDF file and organizes it into a list of dictionaries, 
    each containing the page number and the text extracted from that page.
    
    Args:
    pdf_path (str): The file path of the PDF from which to extract text.
    
    Returns:
    list of dict: A list of dictionaries, each with keys 'page' for the page number and 'content' for the text.
    """
    pages = []
    with open(pdf_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            # Create a dictionary for each page
            page_info = {
                "page": page_num + 1,  # Adding 1 because page_num starts from 0
                "content": page.extract_text() if page.extract_text() else ""  # Adding a check for pages with no text
            }
            pages.append(page_info)

    return pages


def extract_text_from_docx(docx_path):
    """
    Extracts text from a DOCX file using Aspose.Words and organizes it into a list of dictionaries, 
    each containing a page number and the text extracted from that page.
    
    Args:
    docx_path (str): The file path of the DOCX from which to extract text.
    
    Returns:
    list of dict: A list of dictionaries, each with keys 'page' for the page number and 'content' for the text.
    """
    # Load the document
    doc = aw.Document(docx_path)

    # Create a LayoutCollector object which collects layout information
    layout_collector = aw.layout.LayoutCollector(doc)

    pages = []
    current_page = 1
    page_text = []

    # Iterate through each paragraph in the document
    for para in doc.get_child_nodes(aw.NodeType.PARAGRAPH, True):
        node_page = layout_collector.get_start_page_index(para)
        
        # Check if the paragraph is on a new page
        if node_page > current_page:
            pages.append({"page": current_page, "content": " ".join(page_text)})
            page_text = []
            current_page = node_page

        # Extract text from the paragraph
        paragraph_text = para.get_text().strip()
        if paragraph_text:
            page_text.append(paragraph_text)

    # Add the last page
    if page_text:
        pages.append({"page": current_page, "content": " ".join(page_text)})

    return pages





def extract_images_from_docx(docx_path: str, output_folder: str):
    """
    Extracts images from a Word (.docx) document and saves them in the specified output folder.

    Args:
    docx_path (str): Path to the Word document.
    output_folder (str): Folder where extracted images will be saved.
    """
    # Load the Word document
    doc = Document(docx_path)

    # Ensure output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    image_count = 0  # Initialize a counter for images

    # Iterate through each paragraph and table to find images
    for i, element in enumerate(doc.element.body):
        # Check if element is a paragraph
        if element.tag.endswith('p'):
            for run in element.iter():
                if run.tag.endswith('drawing'):
                    # Extract and save the image
                    for inline in run:
                        if inline.tag.endswith('graphic'):
                            image = inline.xpath('.//a:blip')[0].attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}embed']
                            image_part = doc.part.related_parts[image]
                            image_bytes = image_part.blob
                            image_format = image_part.content_type.split('/')[1]

                            # Save the image
                            image_filename = f"paragraph-{i + 1}-img-{image_count + 1}.{image_format}"
                            image_path = os.path.join(output_folder, image_filename)
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_bytes)

                            image_count += 1
        # Check if element is a table
        elif element.tag.endswith('tbl'):
            for row in element.iter():
                if row.tag.endswith('tr'):
                    for cell in row.iter():
                        if cell.tag.endswith('tc'):
                            for paragraph in cell.iter():
                                if paragraph.tag.endswith('p'):
                                    for run in paragraph.iter():
                                        if run.tag.endswith('drawing'):
                                            # Extract and save the image
                                            for inline in run:
                                                if inline.tag.endswith('graphic'):
                                                    image = inline.xpath('.//a:blip')[0].attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}embed']
                                                    image_part = doc.part.related_parts[image]
                                                    image_bytes = image_part.blob
                                                    image_format = image_part.content_type.split('/')[1]

                                                    # Save the image
                                                    image_filename = f"table-{i + 1}-img-{image_count + 1}.{image_format}"
                                                    image_path = os.path.join(output_folder, image_filename)
                                                    with open(image_path, "wb") as img_file:
                                                        img_file.write(image_bytes)

                                                    image_count += 1

    print(f"Extracted {image_count} images from the document.")

